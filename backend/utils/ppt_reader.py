import io
from pptx import Presentation

def extract_text_from_ppt(file_path):
    """
    Extract text from a PPTX file.
    """
    text = ""
    try:
        with open(file_path, 'rb') as f:
            file_bytes = f.read()
            
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        
    return text
